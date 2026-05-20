from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "ppt" / "Shopify_TikTok_AI_Luxury_Resale_OS_SAFE.pptx"

INK = RGBColor(13, 17, 23)
PANEL = RGBColor(22, 27, 34)
PANEL2 = RGBColor(33, 38, 45)
BORDER = RGBColor(48, 54, 61)
TEXT = RGBColor(230, 237, 243)
MUTED = RGBColor(139, 148, 158)
GOLD = RGBColor(200, 164, 93)
GREEN = RGBColor(63, 185, 80)
RED = RGBColor(248, 81, 73)


slides = [
    ("Shopify + TikTok AI Luxury Resale OS", "Shopify controls the asset. TikTok controls the attention. AI controls the workflow.", ["UK + Hong Kong", "Luxury resale + budget fashion separated", "Investor / partner deck"]),
    ("一頁講晒", "呢個唔係普通二手店，而係 AI 一人公司 operating system。", ["Shopify 做中央產品庫、checkout、proof summary、CRM", "TikTok 做短片、直播、廣告、retargeting", "AI Agent Team 做 listing、影片、客服、風控、出貨、數據"]),
    ("Business Model", "四層收入模式", ["Luxury resale cashflow", "Budget fashion for old buyers", "AI live commerce agency", "Future reseller SaaS"]),
    ("UK + HK Market Split", "同一套 OS，兩個市場規則分開。", ["UK: GBP, Shopify UK, TikTok Shop UK, tracked delivery", "HK: HKD, Shopify HK, IG/WhatsApp/Carousell support, SF Express", "客服語言、付款、物流、條款分開"]),
    ("Shopify Main Hub", "Shopify 係資產基地。", ["SKU、價錢、庫存、proof score", "Luxury / budget collection 分開", "UK/HK currency 分開", "弱 proof luxury item 只可以 draft"]),
    ("TikTok Main Growth Engine", "TikTok 係流量同直播成交引擎。", ["每件貨生成 hooks、script、caption、hashtags", "直播前中後都有流程", "TikTok ads + events 做 retargeting", "TikTok Shop API 要等官方 approval"]),
    ("Proof Pack System", "Trust business 要用證據保護。", ["Source record", "Receipt / chat record", "Serial/code 如有", "Detail + flaw photos", "Packing video", "Tracking / payment proof"]),
    ("Proof Score Gate", "Proof 弱就唔可以推 luxury。", ["Strong: publish + promote", "Medium: list carefully", "Weak: manual review / draft", "No proof: block luxury publishing"]),
    ("Video Factory", "一件貨變成 10-20 條內容。", ["Product detail clip", "Close-up / flaw clip", "Styling clip", "Proof/trust clip", "Live preview", "Ad angle"]),
    ("AI Live Commerce", "直播唔係亂講，係流程化。", ["Before live: trailer、VIP preview、proof check", "During live: host script、close-up、Q&A capture", "After live: highlights、CRM follow-up、risk record"]),
    ("Old Buyer + CRM", "舊客係資產，但要分層。", ["Luxury buyer", "Budget buyer", "VIP", "Live viewer", "Risk buyer excluded", "Only contact with marketing consent"]),
    ("Discord + Snapchat", "輔助 channel，唔係主店。", ["Discord: VIP drops、proof education、live reminder", "Snapchat: budget fashion story ads、youth discovery", "流量最後要返 TikTok / Shopify"]),
    ("Fake-Claim Defence", "用 records 回應，唔用情緒爭拗。", ["保存 listing、相、proof pack、packing video、tracking", "Refund / fake claim / chargeback 一律 founder approval", "Return item matching"]),
    ("AI Agent Team", "一個 founder，多個 AI operator。", ["Buyer", "Proof", "Shopify", "TikTok", "Support", "Risk", "Finance", "Live Ops"]),
    ("90-Day Roadmap", "由 prototype 去 repeatable machine。", ["Days 1-30: proof DB + Shopify hub + TikTok content", "Days 31-60: live drops + CRM + event tracking", "Days 61-90: partner pilot + agency offer + dashboard"]),
    ("Next Step", "先用安全版本交付，再接真 API。", ["HTML deck 可直接展示", "SAFE PPTX 可用 PowerPoint 開", "下一步接 Shopify API + OpenAI", "TikTok Shop API 等官方 approval"])
]


def set_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = INK


def add_textbox(slide, left, top, width, height, text, size=24, color=TEXT, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    return box


def add_card(slide, left, top, width, height, title, body="", accent=GOLD):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = BORDER
    bar = slide.shapes.add_shape(1, left, top, Inches(0.07), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    add_textbox(slide, left + Inches(0.22), top + Inches(0.18), width - Inches(0.44), Inches(0.35), title, 17, TEXT, True)
    if body:
        add_textbox(slide, left + Inches(0.22), top + Inches(0.62), width - Inches(0.44), height - Inches(0.75), body, 12, MUTED)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for idx, (title, subtitle, bullets) in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        set_bg(slide)
        add_textbox(slide, Inches(0.55), Inches(0.35), Inches(7.5), Inches(0.3), "SHOPIFY + TIKTOK MAIN", 9, GOLD, True)
        add_textbox(slide, Inches(12.1), Inches(0.35), Inches(0.6), Inches(0.3), f"{idx:02d}", 9, MUTED, True)
        add_textbox(slide, Inches(0.65), Inches(0.9), Inches(8.7), Inches(1.1), title, 34 if idx != 1 else 42, TEXT, True)
        add_textbox(slide, Inches(0.68), Inches(2.0), Inches(8.8), Inches(0.55), subtitle, 17, MUTED)
        if idx == 1:
            x = Inches(0.7)
            for label in bullets:
                add_card(slide, x, Inches(5.55), Inches(3.85), Inches(1.0), label, "", GOLD)
                x += Inches(4.1)
        else:
            cols = 2 if len(bullets) <= 4 else 3
            card_w = Inches(5.75 if cols == 2 else 3.75)
            card_h = Inches(0.95)
            for i, bullet in enumerate(bullets):
                col = i % cols
                row = i // cols
                left = Inches(0.72) + col * (card_w + Inches(0.35))
                top = Inches(3.0) + row * Inches(1.22)
                add_card(slide, left, top, card_w, card_h, bullet, "", GREEN if i % 2 else GOLD)
    prs.save(OUT)
    print(f"Generated {OUT}")


if __name__ == "__main__":
    build_deck()

